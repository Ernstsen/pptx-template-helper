"""Long-form English date rendering and pptx token substitution.

Date rendering (research §R5, FR-009):
    f"{day} {month_name} {year}" with hard-coded English month names.

Token substitution (contracts/template-tokens.md):
    - Date tokens replaced in-place at the paragraph level (surrounding
      text preserved). Paragraph runs are coalesced into one run carrying
      the existing first run's formatting; the documented trade-off is
      that inline (mid-paragraph) run formatting is lost. Paragraph-level
      formatting and the text frame's body properties (`bodyPr`,
      including `normAutofit`) are preserved.
    - Agenda tokens replace the entire text frame: the frame is cleared,
      then one paragraph per issue title is written in given order. The
      first written paragraph reuses the original first run's formatting
      where present so the user's bullet style carries over.
    - Required tokens missing → clear error naming the missing token.
    - {{AGENDA_DEMO}}, {{AGENDA_NO_DEMO}}, or {{AGENDA_OPEN}} appearing
      more than once → explicit error.
"""

from __future__ import annotations

import copy
from datetime import date
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.text.text import _Paragraph, _Run, TextFrame  # noqa: F401  (typing)

from sprint_recap.models import AgendaPlan, AgendaRow, Sprint


_MONTHS = [
    "January", "February", "March", "April", "May", "June",
    "July", "August", "September", "October", "November", "December",
]

DATE_TOKENS = ("{{SPRINT_START}}", "{{SPRINT_END}}", "{{RECAP_DATE}}")
AGENDA_DEMO_TOKEN = "{{AGENDA_DEMO}}"
AGENDA_NO_DEMO_TOKEN = "{{AGENDA_NO_DEMO}}"
AGENDA_OPEN_TOKEN = "{{AGENDA_OPEN}}"
DEMO_START_TAG = "{{DEMO_ITEM_START}}"
DEMO_END_TAG = "{{DEMO_ITEM_END}}"
ITEM_TITLE_TAG = "{{ITEM_TITLE}}"

_LAYOUT_RELTYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
)
_NOTES_RELTYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
)
_R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def format_long_date(d: date) -> str:
    return f"{int(d.day)} {_MONTHS[d.month - 1]} {d.year}"


def _iter_text_frames(prs) -> Iterable:
    """Yield every text frame in slide bodies, including descendants of
    group shapes. Excludes masters, layouts, speaker notes, and chart
    text per template-tokens contract §"Where each token may live"."""
    for slide in prs.slides:
        yield from _iter_shapes_text_frames(slide.shapes)


def _iter_shapes_text_frames(shapes) -> Iterable:
    for shape in shapes:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            yield from _iter_shapes_text_frames(shape.shapes)
            continue
        if shape.has_text_frame:
            yield shape.text_frame


def _replace_in_paragraph(paragraph: _Paragraph, old: str, new: str) -> bool:
    """Replace `old` with `new` inside a paragraph, coalescing runs into the
    first run. Returns True if a replacement was made.

    The paragraph's formatting (font, bullet style) is preserved by reusing
    the first run's `rPr`. Subsequent runs are removed so the paragraph
    holds a single run with the new text.
    """
    full_text = paragraph.text
    if old not in full_text:
        return False
    new_text = full_text.replace(old, new)

    runs = paragraph.runs
    if not runs:
        # Empty paragraph; just set text.
        paragraph.text = new_text
        return True

    # Keep first run's text+formatting, drop the rest.
    first = runs[0]
    first.text = new_text
    # Remove subsequent <a:r> / <a:br> elements while keeping <a:pPr>.
    p_elem = paragraph._p
    keep = first._r
    for child in list(p_elem):
        tag = child.tag.rsplit("}", 1)[-1]
        if tag == "pPr":
            continue
        if child is keep:
            continue
        p_elem.remove(child)
    return True


def _count_token_in_text_frames(text_frames, token: str) -> int:
    n = 0
    for tf in text_frames:
        for para in tf.paragraphs:
            if token in para.text:
                n += para.text.count(token)
    return n


def _clear_text_frame(text_frame) -> _Run | None:
    """Clear all paragraphs from a text frame, preserving the first run's
    rPr template (so the caller can clone its formatting onto subsequent
    paragraphs). Returns the first run before clearing, or None if the
    frame had no runs.

    `python-pptx` does not expose a direct paragraph-delete API, so we
    drop child <a:p> elements from the txBody and let
    `text_frame.text = ""` recreate a single empty paragraph.
    """
    template_run = None
    if text_frame.paragraphs:
        first_runs = text_frame.paragraphs[0].runs
        if first_runs:
            template_run = first_runs[0]

    # Snapshot the first run's rPr (formatting) before we wipe the frame.
    rpr_xml = None
    if template_run is not None and template_run._r.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
    ) is not None:
        rpr_xml = copy.deepcopy(
            template_run._r.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
            )
        )

    text_frame.clear()  # leaves one empty paragraph
    return rpr_xml  # type: ignore[return-value]


def _write_agenda(text_frame, rows: Sequence[AgendaRow]) -> None:
    rpr_xml = _clear_text_frame(text_frame)

    if not rows:
        # Leave the frame empty (one empty paragraph from clear()).
        return

    # First paragraph reuses the cleared frame's existing paragraph object.
    first_para = text_frame.paragraphs[0]
    first_para.text = rows[0].display_title
    if rpr_xml is not None and first_para.runs:
        # Clone the captured rPr onto the new run.
        first_run_elem = first_para.runs[0]._r
        existing_rpr = first_run_elem.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
        )
        if existing_rpr is not None:
            first_run_elem.remove(existing_rpr)
        first_run_elem.insert(0, copy.deepcopy(rpr_xml))

    for row in rows[1:]:
        para = text_frame.add_paragraph()
        para.text = row.display_title
        if rpr_xml is not None and para.runs:
            run_elem = para.runs[0]._r
            existing_rpr = run_elem.find(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
            )
            if existing_rpr is not None:
                run_elem.remove(existing_rpr)
            run_elem.insert(0, copy.deepcopy(rpr_xml))


def _get_slide_notes_text(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text


def _find_item_range(prs) -> tuple[int, int] | None:
    start_idx = end_idx = None
    for i, slide in enumerate(prs.slides):
        notes = _get_slide_notes_text(slide)
        if DEMO_START_TAG in notes:
            if start_idx is not None:
                raise ValueError(
                    f"Multiple {DEMO_START_TAG} tags found in template"
                )
            start_idx = i
        if DEMO_END_TAG in notes:
            if end_idx is not None:
                raise ValueError(
                    f"Multiple {DEMO_END_TAG} tags found in template"
                )
            end_idx = i
    if start_idx is None and end_idx is None:
        return None
    if start_idx is None:
        raise ValueError(
            f"{DEMO_END_TAG} found without matching {DEMO_START_TAG}"
        )
    if end_idx is None:
        raise ValueError(
            f"{DEMO_START_TAG} found without matching {DEMO_END_TAG}"
        )
    if end_idx < start_idx:
        raise ValueError(
            f"{DEMO_END_TAG} appears before {DEMO_START_TAG}"
        )
    return start_idx, end_idx


def _clone_slide(prs, slide):
    new_slide = prs.slides.add_slide(slide.slide_layout)
    rId_map = {}
    for rel in slide.part.rels.values():
        if rel.reltype in (_LAYOUT_RELTYPE, _NOTES_RELTYPE):
            continue
        if rel.is_external:
            new_rId = new_slide.part.rels.get_or_add_ext_rel(
                rel.reltype, rel.target_ref
            )
        else:
            new_rId = new_slide.part.rels.get_or_add(
                rel.reltype, rel.target_part
            )
        rId_map[rel.rId] = new_rId
    new_xml = copy.deepcopy(slide._element)
    for elem in new_xml.iter():
        for attr_name in list(elem.attrib):
            if _R_NS in attr_name:
                old_val = elem.get(attr_name)
                if old_val in rId_map:
                    elem.set(attr_name, rId_map[old_val])
    for child in list(new_slide._element):
        new_slide._element.remove(child)
    for child in list(new_xml):
        new_slide._element.append(child)
    if "shapes" in new_slide.__dict__:
        del new_slide.__dict__["shapes"]
    if slide.has_notes_slide:
        src_body = slide.notes_slide.notes_text_frame._txBody
        dst_body = new_slide.notes_slide.notes_text_frame._txBody
        for child in list(dst_body):
            if child.tag == _A_NS + "p":
                dst_body.remove(child)
        for child in src_body:
            if child.tag == _A_NS + "p":
                dst_body.append(copy.deepcopy(child))
    return new_slide


def _expand_demo_range(prs, demo_rows: Sequence[AgendaRow]) -> None:
    result = _find_item_range(prs)
    if result is None:
        return
    start_idx, end_idx = result
    template_slides = [prs.slides[i] for i in range(start_idx, end_idx + 1)]
    original_count = len(prs.slides)
    for row in demo_rows:
        for tmpl in template_slides:
            new_slide = _clone_slide(prs, tmpl)
            for tf in _iter_shapes_text_frames(new_slide.shapes):
                for para in tf.paragraphs:
                    if ITEM_TITLE_TAG in para.text:
                        _replace_in_paragraph(
                            para, ITEM_TITLE_TAG, row.display_title
                        )
            if new_slide.has_notes_slide:
                for para in new_slide.notes_slide.notes_text_frame.paragraphs:
                    if DEMO_START_TAG in para.text:
                        _replace_in_paragraph(para, DEMO_START_TAG, "")
                    if DEMO_END_TAG in para.text:
                        _replace_in_paragraph(para, DEMO_END_TAG, "")
    sldIdLst = prs.slides._sldIdLst
    all_ids = list(sldIdLst)
    before = all_ids[:start_idx]
    after = all_ids[end_idx + 1:original_count]
    clones = all_ids[original_count:]
    new_order = before + clones + after
    for child in list(sldIdLst):
        sldIdLst.remove(child)
    for sid in new_order:
        sldIdLst.append(sid)


def render_deck(
    template_path: Path,
    output_path: Path,
    sprint: Sprint,
    agenda_plan: AgendaPlan,
) -> None:
    prs = Presentation(str(template_path))

    _expand_demo_range(prs, agenda_plan.demo)

    text_frames = list(_iter_text_frames(prs))

    # --- Validate token presence and cardinality before any edits. ---
    missing: list[str] = []

    def _present(token: str) -> int:
        return _count_token_in_text_frames(text_frames, token)

    sprint_start_n = _present("{{SPRINT_START}}")
    sprint_end_n = _present("{{SPRINT_END}}")
    demo_n = _present(AGENDA_DEMO_TOKEN)
    no_demo_n = _present(AGENDA_NO_DEMO_TOKEN)
    open_n = _present(AGENDA_OPEN_TOKEN)

    if sprint_start_n == 0:
        missing.append("{{SPRINT_START}}")
    if sprint_end_n == 0:
        missing.append("{{SPRINT_END}}")
    if demo_n == 0:
        missing.append(AGENDA_DEMO_TOKEN)
    if no_demo_n == 0:
        missing.append(AGENDA_NO_DEMO_TOKEN)
    if open_n == 0:
        missing.append(AGENDA_OPEN_TOKEN)
    if missing:
        raise ValueError(
            "Required template tokens missing: " + ", ".join(missing)
        )

    if demo_n > 1:
        raise ValueError(
            f"Token {AGENDA_DEMO_TOKEN} appears more than once in the "
            "template; cannot infer which occurrence wins."
        )
    if no_demo_n > 1:
        raise ValueError(
            f"Token {AGENDA_NO_DEMO_TOKEN} appears more than once in the "
            "template; cannot infer which occurrence wins."
        )
    if open_n > 1:
        raise ValueError(
            f"Token {AGENDA_OPEN_TOKEN} appears more than once in the "
            "template; cannot infer which occurrence wins."
        )

    # --- Date substitution (paragraph-level, in place). ---
    start_str = format_long_date(sprint.start)
    end_str = format_long_date(sprint.end)
    recap_str = end_str  # FR-009: recap-meeting date == sprint end.

    for tf in text_frames:
        for para in tf.paragraphs:
            text = para.text
            if (
                "{{SPRINT_START}}" in text
                or "{{SPRINT_END}}" in text
                or "{{RECAP_DATE}}" in text
                or "{{SPRINT_ID}}" in text
            ):
                if "{{SPRINT_START}}" in para.text:
                    _replace_in_paragraph(para, "{{SPRINT_START}}", start_str)
                if "{{SPRINT_END}}" in para.text:
                    _replace_in_paragraph(para, "{{SPRINT_END}}", end_str)
                if "{{RECAP_DATE}}" in para.text:
                    _replace_in_paragraph(para, "{{RECAP_DATE}}", recap_str)
                if "{{SPRINT_ID}}" in para.text:
                    _replace_in_paragraph(para, "{{SPRINT_ID}}", sprint.name)

    # --- Agenda substitution (clear + write per issue). ---
    for tf in text_frames:
        joined = "\n".join(p.text for p in tf.paragraphs)
        if AGENDA_DEMO_TOKEN in joined:
            _write_agenda(tf, agenda_plan.demo)
        elif AGENDA_NO_DEMO_TOKEN in joined:
            _write_agenda(tf, agenda_plan.no_demo)
        elif AGENDA_OPEN_TOKEN in joined:
            _write_agenda(tf, agenda_plan.open)

    prs.save(str(output_path))
