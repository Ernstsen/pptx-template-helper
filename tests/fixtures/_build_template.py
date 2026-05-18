"""Build the minimal fixture template used by tests/unit/test_deck_tokens.py.

The committed `tests/fixtures/template.pptx` should be the output of running
this script with python-pptx 1.0.2 (the pinned version in
`requirements.txt`). Re-run after editing this file to keep the fixture
in sync:

    python tests/fixtures/_build_template.py

The fixture deliberately includes:

  - one slide with a title containing both `{{SPRINT_START}}` and
    `{{SPRINT_END}}` and an extra prose suffix so paragraph-level
    substitution is exercised with surrounding text (template-tokens
    contract §"Date tokens");
  - a subtitle text frame containing `{{RECAP_DATE}}`;
  - a footer text frame containing `{{SPRINT_ID}}`;
  - a second slide with three text frames whose sole content is
    `{{AGENDA_DEMO}}`, `{{AGENDA_NO_DEMO}}`, and `{{AGENDA_OPEN}}`
    respectively, matching the agenda-token contract.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def build(out_path: Path) -> None:
    prs = Presentation()

    # Slide 1: title slide with date tokens.
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    slide1.shapes.title.text = "Sprint {{SPRINT_START}} – {{SPRINT_END}}"
    # Subtitle placeholder index varies by template; layout 0 has it at idx 1.
    slide1.placeholders[1].text = "Recap meeting: {{RECAP_DATE}}"

    footer = slide1.shapes.add_textbox(
        Inches(0.3), Inches(6.5), Inches(4.0), Inches(0.5)
    )
    footer.text_frame.text = "{{SPRINT_ID}}"

    # Slide 2: blank layout with three text-box agenda regions.
    blank = prs.slide_layouts[6]
    slide2 = prs.slides.add_slide(blank)

    demo_box = slide2.shapes.add_textbox(
        Inches(0.3), Inches(0.5), Inches(3.0), Inches(6)
    )
    demo_box.text_frame.text = "{{AGENDA_DEMO}}"

    no_demo_box = slide2.shapes.add_textbox(
        Inches(3.5), Inches(0.5), Inches(3.0), Inches(6)
    )
    no_demo_box.text_frame.text = "{{AGENDA_NO_DEMO}}"

    open_box = slide2.shapes.add_textbox(
        Inches(6.7), Inches(0.5), Inches(3.0), Inches(6)
    )
    open_box.text_frame.text = "{{AGENDA_OPEN}}"

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


if __name__ == "__main__":
    target = Path(__file__).parent / "template.pptx"
    build(target)
    print(f"wrote {target}")
