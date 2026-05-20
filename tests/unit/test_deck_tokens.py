"""Token substitution against tests/fixtures/template.pptx (FR-009/FR-010,
contracts/template-tokens.md)."""

from __future__ import annotations

import shutil
from datetime import date, datetime, timezone
from pathlib import Path

import pytest
from pptx import Presentation

from sprint_recap.deck import render_deck
from sprint_recap.models import AgendaPlan, AgendaRow, Sprint, SprintIssue


FIXTURE = Path(__file__).resolve().parents[1] / "fixtures" / "template.pptx"


def _sprint() -> Sprint:
    return Sprint(
        id="121-318",
        name="Sprint 42",
        start=date(2026, 4, 8),
        end=date(2026, 5, 5),
        archived=True,
    )


def _issue(id_: str, title: str, resolved: bool) -> SprintIssue:
    return SprintIssue(
        id_readable=id_,
        title=title,
        issue_type="Story",
        parent_id_readable=None,
        resolved_at=datetime(2026, 4, 30, tzinfo=timezone.utc) if resolved else None,
        created_at=datetime(2026, 4, 9, tzinfo=timezone.utc),
    )


def _agenda() -> AgendaPlan:
    return AgendaPlan(
        demo=[
            AgendaRow.from_issue(
                _issue("PROJ-1", "Migrate billing service to v3", resolved=True)
            ),
        ],
        no_demo=[
            AgendaRow.from_issue(
                _issue("PROJ-2", "Hot-patch invoice rounding", resolved=True)
            ),
        ],
        open=[
            AgendaRow.from_issue(
                _issue(
                    "PROJ-3",
                    "Investigate p99 latency on report API",
                    resolved=False,
                )
            ),
        ],
        unfiltered_count=3,
        filtered_count=3,
        collapsed_subtask_count=0,
    )


def _all_paragraph_text(out_path: Path) -> list[str]:
    prs = Presentation(str(out_path))
    out: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    out.append(para.text)
    return out


def test_render_deck_substitutes_date_tokens_in_place(tmp_path: Path) -> None:
    template = tmp_path / "Recap-Template.pptx"
    shutil.copy(FIXTURE, template)
    output = tmp_path / "out.pptx"

    render_deck(template, output, _sprint(), _agenda())

    paragraphs = _all_paragraph_text(output)
    # Date tokens replaced; surrounding text preserved.
    assert "Sprint 8 April 2026 – 5 May 2026" in paragraphs
    assert "Recap meeting: 5 May 2026" in paragraphs
    # No raw token left behind.
    joined = "\n".join(paragraphs)
    assert "{{SPRINT_START}}" not in joined
    assert "{{SPRINT_END}}" not in joined
    assert "{{RECAP_DATE}}" not in joined
    assert "{{SPRINT_ID}}" not in joined
    assert "{{AGENDA_DEMO}}" not in joined
    assert "{{AGENDA_NO_DEMO}}" not in joined
    assert "{{AGENDA_OPEN}}" not in joined


def test_render_deck_substitutes_sprint_id_token(tmp_path: Path) -> None:
    template = tmp_path / "Recap-Template.pptx"
    shutil.copy(FIXTURE, template)
    output = tmp_path / "out.pptx"

    render_deck(template, output, _sprint(), _agenda())

    paragraphs = _all_paragraph_text(output)
    assert "Sprint 42" in paragraphs


def test_render_deck_writes_one_paragraph_per_issue_in_order(tmp_path: Path) -> None:
    template = tmp_path / "Recap-Template.pptx"
    shutil.copy(FIXTURE, template)
    output = tmp_path / "out.pptx"

    render_deck(template, output, _sprint(), _agenda())

    prs = Presentation(str(output))
    # Slide 2 has three agenda boxes (per fixture _build_template.py).
    slide2 = prs.slides[1]
    demo_box = slide2.shapes[0]
    no_demo_box = slide2.shapes[1]
    open_box = slide2.shapes[2]

    demo_paragraphs = [p.text for p in demo_box.text_frame.paragraphs]
    no_demo_paragraphs = [p.text for p in no_demo_box.text_frame.paragraphs]
    open_paragraphs = [p.text for p in open_box.text_frame.paragraphs]

    assert demo_paragraphs == ["Migrate billing service to v3"]
    assert no_demo_paragraphs == ["Hot-patch invoice rounding"]
    assert open_paragraphs == ["Investigate p99 latency on report API"]


def test_render_deck_does_not_modify_template_bytes(tmp_path: Path) -> None:
    template = tmp_path / "Recap-Template.pptx"
    shutil.copy(FIXTURE, template)
    output = tmp_path / "out.pptx"

    before = template.read_bytes()
    render_deck(template, output, _sprint(), _agenda())
    after = template.read_bytes()

    assert before == after, "template was modified in place"
    assert output.exists()


def test_render_deck_raises_when_required_token_missing(tmp_path: Path) -> None:
    """If a required token is removed from the template, the program names
    the missing token in its error (template-tokens contract / spec edge case)."""
    template = tmp_path / "Recap-Template.pptx"
    shutil.copy(FIXTURE, template)

    # Strip {{AGENDA_DEMO}} from the copy by replacing it with empty text.
    prs = Presentation(str(template))
    slide2 = prs.slides[1]
    slide2.shapes[0].text_frame.text = "(removed)"
    prs.save(str(template))

    output = tmp_path / "out.pptx"
    with pytest.raises(Exception) as exc:
        render_deck(template, output, _sprint(), _agenda())
    assert "{{AGENDA_DEMO}}" in str(exc.value)


def test_render_deck_writes_display_title_not_issue_title(tmp_path: Path) -> None:
    """Spec 005: the renderer reads `display_title`, so an edited row
    shows the user's text instead of the original YouTrack title."""
    template = tmp_path / "Recap-Template.pptx"
    shutil.copy(FIXTURE, template)
    output = tmp_path / "out.pptx"

    plan = _agenda()
    plan.no_demo[0].display_title = "Edited mention title"
    plan.open[0].display_title = "Edited open title"

    render_deck(template, output, _sprint(), plan)
    paragraphs = _all_paragraph_text(output)
    assert "Edited mention title" in paragraphs
    assert "Edited open title" in paragraphs
    # Original (un-edited) titles must not have been written for the
    # rows we renamed.
    assert "Hot-patch invoice rounding" not in paragraphs
    assert "Investigate p99 latency on report API" not in paragraphs


def test_render_deck_raises_on_duplicate_agenda_token(tmp_path: Path) -> None:
    """Duplicate {{AGENDA_DEMO}} → explicit error (template-tokens contract)."""
    from pptx.util import Inches

    template = tmp_path / "Recap-Template.pptx"
    shutil.copy(FIXTURE, template)

    # Add an extra text box also containing {{AGENDA_DEMO}} so the
    # required-token-presence check passes but the cardinality check fails.
    prs = Presentation(str(template))
    slide2 = prs.slides[1]
    extra = slide2.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(1))
    extra.text_frame.text = "{{AGENDA_DEMO}}"
    prs.save(str(template))

    output = tmp_path / "out.pptx"
    with pytest.raises(Exception) as exc:
        render_deck(template, output, _sprint(), _agenda())
    msg = str(exc.value)
    assert "{{AGENDA_DEMO}}" in msg and (
        "duplicate" in msg.lower() or "more than once" in msg.lower()
    )
