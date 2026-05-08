"""Demo item slide range expansion (spec 003-demo-item-slides)."""

from __future__ import annotations

import shutil
import struct
import zlib
from datetime import date, datetime, timezone
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from sprint_recap.deck import render_deck
from sprint_recap.models import AgendaPlan, Sprint, SprintIssue

FIXTURE = Path(__file__).resolve().parents[1] / "fixtures" / "template.pptx"


def _sprint() -> Sprint:
    return Sprint(
        id="121-318",
        name="Sprint 42",
        start=date(2026, 4, 8),
        end=date(2026, 5, 5),
        archived=True,
    )


def _issue(id_: str, title: str, resolved: bool = True) -> SprintIssue:
    return SprintIssue(
        id_readable=id_,
        title=title,
        issue_type="Story",
        parent_id_readable=None,
        resolved_at=datetime(2026, 4, 30, tzinfo=timezone.utc) if resolved else None,
        created_at=datetime(2026, 4, 9, tzinfo=timezone.utc),
    )


def _agenda(
    demo: list[SprintIssue] | None = None,
    no_demo: list[SprintIssue] | None = None,
    open_: list[SprintIssue] | None = None,
) -> AgendaPlan:
    return AgendaPlan(
        demo=demo if demo is not None else [],
        no_demo=no_demo if no_demo is not None else [_issue("P-10", "No-demo item")],
        open=open_ if open_ is not None else [_issue("P-11", "Open item", resolved=False)],
        unfiltered_count=0,
        filtered_count=0,
    )


def _add_required_slides(prs: Presentation) -> None:
    """Add slides with all required tokens (date + agenda)."""
    title_layout = prs.slide_layouts[0]
    s1 = prs.slides.add_slide(title_layout)
    s1.shapes.title.text = "Sprint {{SPRINT_START}} – {{SPRINT_END}}"
    s1.placeholders[1].text = "Recap meeting: {{RECAP_DATE}}"

    blank = prs.slide_layouts[6]
    s2 = prs.slides.add_slide(blank)
    for token, x in [
        ("{{AGENDA_DEMO}}", 0.3),
        ("{{AGENDA_NO_DEMO}}", 3.5),
        ("{{AGENDA_OPEN}}", 6.7),
    ]:
        box = s2.shapes.add_textbox(
            Inches(x), Inches(0.5), Inches(3), Inches(6)
        )
        box.text_frame.text = token


def _build_demo_template(path: Path, *, range_size: int = 2) -> None:
    prs = Presentation()
    _add_required_slides(prs)

    blank = prs.slide_layouts[6]
    for i in range(range_size):
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(
            Inches(1), Inches(1), Inches(8), Inches(2)
        )
        tb.text_frame.text = "Demo: {{ITEM_TITLE}}"
        parts = []
        if i == 0:
            parts.append("{{DEMO_ITEM_START}}")
        if i == range_size - 1:
            parts.append("{{DEMO_ITEM_END}}")
        if parts:
            s.notes_slide.notes_text_frame.text = " ".join(parts)

    prs.save(str(path))


def _make_tiny_png(path: Path) -> None:
    width, height = 1, 1
    raw_data = b"\x00\xff\x00\x00"

    def chunk(chunk_type: bytes, data: bytes) -> bytes:
        c = chunk_type + data
        crc = struct.pack(">I", zlib.crc32(c) & 0xFFFFFFFF)
        return struct.pack(">I", len(data)) + c + crc

    ihdr_data = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    with open(path, "wb") as f:
        f.write(b"\x89PNG\r\n\x1a\n")
        f.write(chunk(b"IHDR", ihdr_data))
        f.write(chunk(b"IDAT", zlib.compress(raw_data)))
        f.write(chunk(b"IEND", b""))


# --- Tests ---


def test_two_slide_range_three_issues(tmp_path: Path) -> None:
    template = tmp_path / "tmpl.pptx"
    _build_demo_template(template, range_size=2)
    output = tmp_path / "out.pptx"

    agenda = _agenda(
        demo=[
            _issue("P-1", "Feature Alpha"),
            _issue("P-2", "Feature Beta"),
            _issue("P-3", "Feature Gamma"),
        ],
    )
    render_deck(template, output, _sprint(), agenda)

    prs = Presentation(str(output))
    assert len(prs.slides) == 8

    demo_texts = []
    for i in range(2, 8):
        slide = prs.slides[i]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    demo_texts.append(p.text)

    assert demo_texts.count("Demo: Feature Alpha") == 2
    assert demo_texts.count("Demo: Feature Beta") == 2
    assert demo_texts.count("Demo: Feature Gamma") == 2

    alpha_idx = demo_texts.index("Demo: Feature Alpha")
    beta_idx = demo_texts.index("Demo: Feature Beta")
    gamma_idx = demo_texts.index("Demo: Feature Gamma")
    assert alpha_idx < beta_idx < gamma_idx

    for i in range(2, 8):
        slide = prs.slides[i]
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            assert "{{DEMO_ITEM_START}}" not in notes
            assert "{{DEMO_ITEM_END}}" not in notes


def test_no_range_tags_renders_identically(tmp_path: Path) -> None:
    template = tmp_path / "tmpl.pptx"
    shutil.copy(FIXTURE, template)
    output = tmp_path / "out.pptx"

    agenda = _agenda(
        demo=[_issue("P-1", "Demo item")],
        no_demo=[_issue("P-2", "No-demo item")],
        open_=[_issue("P-3", "Open item", resolved=False)],
    )
    render_deck(template, output, _sprint(), agenda)

    prs = Presentation(str(output))
    assert len(prs.slides) == 2


def test_missing_end_tag_raises(tmp_path: Path) -> None:
    prs = Presentation()
    _add_required_slides(prs)
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    s.notes_slide.notes_text_frame.text = "{{DEMO_ITEM_START}}"
    path = tmp_path / "tmpl.pptx"
    prs.save(str(path))

    with pytest.raises(ValueError, match="DEMO_ITEM_END"):
        render_deck(path, tmp_path / "out.pptx", _sprint(), _agenda())


def test_missing_start_tag_raises(tmp_path: Path) -> None:
    prs = Presentation()
    _add_required_slides(prs)
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    s.notes_slide.notes_text_frame.text = "{{DEMO_ITEM_END}}"
    path = tmp_path / "tmpl.pptx"
    prs.save(str(path))

    with pytest.raises(ValueError, match="DEMO_ITEM_START"):
        render_deck(path, tmp_path / "out.pptx", _sprint(), _agenda())


def test_end_before_start_raises(tmp_path: Path) -> None:
    prs = Presentation()
    _add_required_slides(prs)
    blank = prs.slide_layouts[6]
    s_end = prs.slides.add_slide(blank)
    s_end.notes_slide.notes_text_frame.text = "{{DEMO_ITEM_END}}"
    s_start = prs.slides.add_slide(blank)
    s_start.notes_slide.notes_text_frame.text = "{{DEMO_ITEM_START}}"
    path = tmp_path / "tmpl.pptx"
    prs.save(str(path))

    with pytest.raises(ValueError, match="before"):
        render_deck(path, tmp_path / "out.pptx", _sprint(), _agenda())


def test_empty_demo_removes_range(tmp_path: Path) -> None:
    template = tmp_path / "tmpl.pptx"
    _build_demo_template(template, range_size=2)
    output = tmp_path / "out.pptx"

    render_deck(template, output, _sprint(), _agenda(demo=[]))

    prs = Presentation(str(output))
    assert len(prs.slides) == 2


def test_single_slide_range(tmp_path: Path) -> None:
    template = tmp_path / "tmpl.pptx"
    _build_demo_template(template, range_size=1)
    output = tmp_path / "out.pptx"

    agenda = _agenda(demo=[_issue("P-1", "Solo Feature")])
    render_deck(template, output, _sprint(), agenda)

    prs = Presentation(str(output))
    assert len(prs.slides) == 3

    slide = prs.slides[2]
    texts = [
        p.text
        for s in slide.shapes
        if s.has_text_frame
        for p in s.text_frame.paragraphs
    ]
    assert "Demo: Solo Feature" in texts


def test_item_title_outside_range_unchanged(tmp_path: Path) -> None:
    prs = Presentation()
    _add_required_slides(prs)

    blank = prs.slide_layouts[6]
    s_demo = prs.slides.add_slide(blank)
    tb = s_demo.shapes.add_textbox(
        Inches(1), Inches(1), Inches(8), Inches(2)
    )
    tb.text_frame.text = "Demo: {{ITEM_TITLE}}"
    s_demo.notes_slide.notes_text_frame.text = (
        "{{DEMO_ITEM_START}} {{DEMO_ITEM_END}}"
    )

    s_extra = prs.slides.add_slide(blank)
    tb2 = s_extra.shapes.add_textbox(
        Inches(1), Inches(1), Inches(5), Inches(2)
    )
    tb2.text_frame.text = "Not replaced: {{ITEM_TITLE}}"

    path = tmp_path / "tmpl.pptx"
    prs.save(str(path))
    output = tmp_path / "out.pptx"

    agenda = _agenda(demo=[_issue("P-1", "Some Feature")])
    render_deck(path, output, _sprint(), agenda)

    prs_out = Presentation(str(output))
    assert len(prs_out.slides) == 4
    last = prs_out.slides[3]
    texts = [
        p.text
        for s in last.shapes
        if s.has_text_frame
        for p in s.text_frame.paragraphs
    ]
    assert "Not replaced: {{ITEM_TITLE}}" in texts


def test_notes_text_preserved_around_tags(tmp_path: Path) -> None:
    prs = Presentation()
    _add_required_slides(prs)

    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    tb = s.shapes.add_textbox(
        Inches(1), Inches(1), Inches(8), Inches(2)
    )
    tb.text_frame.text = "{{ITEM_TITLE}}"
    s.notes_slide.notes_text_frame.text = (
        "Presenter notes {{DEMO_ITEM_START}} {{DEMO_ITEM_END}} end"
    )

    path = tmp_path / "tmpl.pptx"
    prs.save(str(path))
    output = tmp_path / "out.pptx"

    agenda = _agenda(demo=[_issue("P-1", "Feature X")])
    render_deck(path, output, _sprint(), agenda)

    prs_out = Presentation(str(output))
    slide = prs_out.slides[2]
    notes = slide.notes_slide.notes_text_frame.text
    assert "Presenter notes" in notes
    assert "end" in notes
    assert "{{DEMO_ITEM_START}}" not in notes
    assert "{{DEMO_ITEM_END}}" not in notes


def test_clone_preserves_images(tmp_path: Path) -> None:
    img_path = tmp_path / "test.png"
    _make_tiny_png(img_path)

    prs = Presentation()
    _add_required_slides(prs)

    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    s.shapes.add_picture(str(img_path), Inches(1), Inches(3), Inches(2), Inches(2))
    tb = s.shapes.add_textbox(
        Inches(1), Inches(1), Inches(8), Inches(1)
    )
    tb.text_frame.text = "{{ITEM_TITLE}}"
    s.notes_slide.notes_text_frame.text = (
        "{{DEMO_ITEM_START}} {{DEMO_ITEM_END}}"
    )

    path = tmp_path / "tmpl.pptx"
    prs.save(str(path))
    output = tmp_path / "out.pptx"

    agenda = _agenda(
        demo=[_issue("P-1", "With Image"), _issue("P-2", "Also Image")],
    )
    render_deck(path, output, _sprint(), agenda)

    prs_out = Presentation(str(output))
    assert len(prs_out.slides) == 4

    for i in [2, 3]:
        slide = prs_out.slides[i]
        has_picture = any(
            shape.shape_type == 13 for shape in slide.shapes
        )
        assert has_picture, f"Slide {i} missing picture"
